import json
import logging
import shutil
import subprocess
from pathlib import Path

import cv2
from pptx import Presentation
from pptx.util import Pt
from wand.image import Image as WandImage

from llm_integration import LLMFactory
from utils import ppt_prompt


def _convert_ppt_to_images(ppt_path: str, output_dir: str):
    """将PPT转换为图片序列"""
    try:
        # 使用 unoserver 进行转换
        cmd = f'unoconvert --port 2002 "{ppt_path}" "{output_dir}/slides.pdf"'
        subprocess.run(cmd, shell=True, check=True)

        # 使用 wand 将 PDF 转换为图片
        with WandImage(filename=f"{output_dir}/slides.pdf", resolution=300) as img:
            img.compression_quality = 95  # 设置图片质量
            img.background_color = "white"  # 设置背景为白色
            img.alpha_channel = 'remove'  # 移除透明通道
            img.save(filename=f"{output_dir}/slide_%03d.png")

    except Exception as e:
        logging.error(f"PPT转换图片失败: {str(e)}")
        raise


def _check_dependencies():
    """检查系统是否安装了必要的依赖工具"""
    required_tools = ['unoconvert', 'ffmpeg']  # 移除 'magick'
    missing_tools = []
    for tool in required_tools:
        if not shutil.which(tool):
            missing_tools.append(tool)
    if missing_tools:
        raise RuntimeError(f"缺少必要的依赖工具: {', '.join(missing_tools)}。请安装后再试。")


class MultimediaGenerator:
    """多媒体资源生成器"""

    def __init__(self, tongyi_api_key: str):
        _check_dependencies()  # 初始化时检查依赖
        self.llm = LLMFactory.initialize_tongyi(api_key=tongyi_api_key)
        self.output_dir = Path("./output/multimedia")
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_ppt_with_narration(self, content: str) -> tuple[str, list]:
        """生成PPT和配音文案"""
        try:
            # 调用LLM生成PPT结构和配音文案
            prompt = ppt_prompt.format(content=content)
            response = self.llm.invoke(prompt)
            logging.debug(f"LLM Response: {response}")

            # 解析LLM返回的JSON
            ppt_data = json.loads(response.strip())
            if not ppt_data or "title" not in ppt_data or "sections" not in ppt_data:
                raise ValueError("LLM返回的PPT结构无效")

            # 创建PPT
            prs = Presentation()
            narrations = []

            # 添加封面
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = ppt_data["title"]
            narrations.append(ppt_data.get("title_narration", "默认封面讲解"))

            # 添加内容页
            for section in ppt_data["sections"]:
                content_slide = prs.slides.add_slide(prs.slide_layouts[1])
                content_slide.shapes.title.text = section.get("title", "默认章节标题")
                content_shape = content_slide.shapes.placeholders[1]
                tf = content_shape.text_frame
                for point in section.get("points", ["默认内容"]):
                    p = tf.add_paragraph()
                    p.text = point
                    p.font.size = Pt(18)
                narrations.append(section.get("narration", "默认讲解文案"))

            # 保存PPT
            ppt_path = self.output_dir / "teaching.pptx"
            prs.save(str(ppt_path))
            return str(ppt_path), narrations

        except Exception as e:
            logging.error(f"PPT生成失败: {str(e)}")
            return "", []

    def _generate_audio(self, text: str, index: int) -> str:
        """生成语音文件"""
        try:
            audio_path = str(self.output_dir / f"narration_{index}.mp3")
            cmd = f'echo "{text}" | text2wave -o "{audio_path}"'
            subprocess.run(cmd, shell=True, check=True)
            return audio_path
        except Exception as e:
            logging.error(f"音频生成失败: {str(e)}")
            return ""

    def _merge_video_audio(self, video_path: str, audio_paths: list, output_path: str):
        """合并视频和音频"""
        try:
            # 合并所有音频文件
            audio_list_file = self.output_dir / "audio_list.txt"
            with open(audio_list_file, 'w') as f:
                for audio_path in audio_paths:
                    f.write(f"file '{audio_path}'\n")

            merged_audio = self.output_dir / "merged_audio.mp3"
            cmd = f'ffmpeg -f concat -safe 0 -i "{audio_list_file}" -c copy "{merged_audio}"'
            subprocess.run(cmd, shell=True, check=True)

            # 合并视频和音频
            cmd = f'ffmpeg -i "{video_path}" -i "{merged_audio}" -c:v copy -c:a aac "{output_path}"'
            subprocess.run(cmd, shell=True, check=True)

        except Exception as e:
            logging.error(f"音视频合并失败: {str(e)}")
            raise

    def generate_teaching_video(self, content: str) -> str:
        """生成教学视频"""
        try:
            _check_dependencies()  # 检查依赖
            # 生成PPT和配音文案
            ppt_path, narrations = self.generate_ppt_with_narration(content)
            if not ppt_path:
                raise Exception("PPT生成失败")

            # 生成语音文件
            audio_paths = []
            for i, narration in enumerate(narrations):
                audio_path = self._generate_audio(narration, i)
                if audio_path:
                    audio_paths.append(audio_path)

            # 将PPT转换为图片序列
            images_dir = self.output_dir / "slides"
            images_dir.mkdir(exist_ok=True)
            _convert_ppt_to_images(ppt_path, str(images_dir))

            # 生成视频
            video_path = self._create_video_from_images(images_dir)
            if not video_path:
                raise Exception("视频生成失败")

            # 合并视频和音频
            final_video_path = str(self.output_dir / "teaching_video.mp4")
            self._merge_video_audio(video_path, audio_paths, final_video_path)

            return final_video_path

        except RuntimeError as e:
            logging.error(f"依赖工具缺失: {str(e)}")
            return ""
        except Exception as e:
            logging.error(f"视频生成失败: {str(e)}")
            return ""

    def _create_video_from_images(self, images_dir: Path) -> str:
        """从图片序列生成视频"""
        try:
            image_files = sorted([f for f in images_dir.glob("slide_*.png")])
            if not image_files:
                raise Exception("未找到幻灯片图片")

            # 创建视频写入器
            first_img = cv2.imread(str(image_files[0]))
            height, width = first_img.shape[:2]
            video_path = str(self.output_dir / "teaching_video_temp.mp4")
            fourcc = cv2.VideoWriter.fourcc(*'H264')
            fps = 30
            out = cv2.VideoWriter(video_path, fourcc, fps, (width, height))

            # 将图片写入视频
            for img_path in image_files:
                img = cv2.imread(str(img_path))
                for _ in range(5):  # 每张图片持续5秒
                    out.write(img)

            out.release()
            return video_path

        except Exception as e:
            logging.error(f"视频生成失败: {str(e)}")
            return ""


# 使用示例
if __name__ == "__main__":
    generator = MultimediaGenerator(tongyi_api_key="sk-7db9f91fb10e4319b4cc04f763c0e461")

    content = """
    主题：勾股定理
    教学目标：理解并掌握勾股定理的概念和应用
    重点内容：
    1. 直角三角形的性质
    2. 勾股定理公式（a² + b² = c²）
    3. 实际应用场景（如测量距离）
    """

    video_path = generator.generate_teaching_video(content)
    print(f"生成的视频路径: {video_path}")