from pptx import Presentation
import re
from src.PptTurnVideo.Polish_Outline import text_polishing
from src.configs import config


def extract_speaker_notes(file_path):
    """提取PPTX中的演讲者备注内容，返回每页备注的列表"""
    prs = Presentation(file_path)
    notes_list = []

    # 遍历所有幻灯片
    for slide in prs.slides:
        page_notes = []
        notes_slide = slide.notes_slide

        # 提取当前页备注内容
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                # 合并段落文本
                page_text = ''.join([paragraph.text for paragraph in shape.text_frame.paragraphs])
                page_notes.append(page_text)

        # 处理合并后的备注文本
        combined_notes = ''.join(page_notes)
        # 使用正则表达式去除末尾数字
        cleaned_notes = re.sub(r'\d+$', '', combined_notes)
        notes_list.append(cleaned_notes)
    api_key = config['api_keys']['deepseek_api_key']
    polish_texts = text_polishing(api_key, notes_list)

    return polish_texts
    # return notes_list

# notes = extract_speaker_notes("downloaded_file.pptx")
# # print(f"演讲者备注内容：\n{notes}")
# print(notes)